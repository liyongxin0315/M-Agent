"""
PPT Generation Node - PPT 生成节点

集成 ppt-generation 技能，提供 PPT 生成能力。
"""

import logging
from typing import Any, Dict, List, Optional, TYPE_CHECKING
from pathlib import Path

from ..base_node import BaseNode, NodeResult, NodeStatus

if TYPE_CHECKING:
    from pptx import Presentation

logger = logging.getLogger(__name__)


class PPTGenerationNode(BaseNode):
    """PPT 生成节点"""
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        super().__init__("ppt_generation", config)
        self._default_theme = config.get("theme", "default") if config else "default"
    
    async def execute(self, context: Dict[str, Any]) -> NodeResult:
        """
        执行 PPT 生成
        
        Args:
            context: 执行上下文，包含:
                - title: PPT 标题
                - content: 内容大纲或文本
                - theme: 主题
                - output_path: 输出路径
                - slides: 幻灯片列表 (可选)
        
        Returns:
            NodeResult: 生成的 PPT 信息
        """
        try:
            title = context.get("title")
            content = context.get("content")
            theme = context.get("theme", self._default_theme)
            output_path = context.get("output_path")
            slides = context.get("slides", [])
            
            if not title and not slides:
                return NodeResult(
                    status=NodeStatus.FAILED,
                    error="缺少必需参数：title 或 slides",
                    node_name=self.name
                )
            
            # 调用 ppt-generation 技能
            result = await self._generate_ppt(
                title=title,
                content=content,
                theme=theme,
                slides=slides,
                output_path=output_path
            )
            
            return NodeResult(
                status=NodeStatus.COMPLETED,
                output=result,
                node_name=self.name
            )
        
        except Exception as e:
            logger.error(f"PPT 生成失败：{e}")
            return NodeResult(
                status=NodeStatus.FAILED,
                error=str(e),
                node_name=self.name
            )
    
    async def _generate_ppt(
        self,
        title: Optional[str] = None,
        content: Optional[str] = None,
        theme: str = "default",
        slides: Optional[List[Dict[str, Any]]] = None,
        output_path: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        生成 PPT
        """
        from pathlib import Path
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        if not output_path:
            output_path = "presentation.pptx"
        
        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        
        # 创建演示文稿
        prs = Presentation()
        
        # 添加标题页
        if title:
            slide_layout = prs.slide_layouts[0]  # Title Slide
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            
            if content:
                # 解析内容为大纲
                outline_slides = self._parse_content_to_slides(content)
                for slide_content in outline_slides:
                    self._add_content_slide(prs, slide_content)
        
        # 添加自定义幻灯片
        if slides:
            for slide_data in slides:
                self._add_custom_slide(prs, slide_data)
        
        # 保存
        prs.save(str(path))
        
        return {
            "output_path": str(path),
            "slides_count": len(prs.slides),
            "theme": theme,
            "title": title
        }
    
    def _parse_content_to_slides(self, content: str) -> List[Dict[str, Any]]:
        """
        解析内容为幻灯片列表
        """
        slides = []
        current_slide = {"title": "", "content": []}
        
        for line in content.split("\n"):
            line = line.strip()
            if not line:
                continue
            
            if line.startswith("#"):
                # 新幻灯片标题
                if current_slide["title"]:
                    slides.append(current_slide)
                current_slide = {"title": line.lstrip("#").strip(), "content": []}
            else:
                # 内容
                current_slide["content"].append(line)
        
        if current_slide["title"]:
            slides.append(current_slide)
        
        return slides
    
    def _add_content_slide(self, prs: "Presentation", slide_data: Dict[str, Any]) -> None:
        """
        添加内容幻灯片
        """
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        slide.shapes.title.text = slide_data["title"]
        
        # 设置内容
        if slide_data["content"]:
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.text = slide_data["content"][0]
            
            for item in slide_data["content"][1:]:
                p = text_frame.add_paragraph()
                p.text = item
                p.level = 0
    
    def _add_custom_slide(self, prs: "Presentation", slide_data: Dict[str, Any]) -> None:
        """
        添加自定义幻灯片
        """
        layout_index = slide_data.get("layout", 1)
        slide_layout = prs.slide_layouts[layout_index] if layout_index < len(prs.slide_layouts) else prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        if "title" in slide_data:
            slide.shapes.title.text = slide_data["title"]
        
        if "content" in slide_data:
            content_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if content_placeholder:
                text_frame = content_placeholder.text_frame
                text_frame.text = slide_data["content"]
    
    def get_schema(self) -> Dict[str, Any]:
        """返回节点输入输出 schema"""
        return {
            "inputs": {
                "title": {"type": "string", "required": False, "description": "PPT 标题"},
                "content": {"type": "string", "required": False, "description": "内容大纲或文本"},
                "theme": {"type": "string", "required": False, "default": "default"},
                "output_path": {"type": "string", "required": False, "description": "输出文件路径"},
                "slides": {
                    "type": "array",
                    "items": {"type": "object"},
                    "required": False,
                    "description": "自定义幻灯片列表"
                }
            },
            "outputs": {
                "generated_ppt": {"type": "object", "description": "生成的 PPT 信息"}
            }
        }
